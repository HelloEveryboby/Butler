from pptx import Presentation

def convert_pptx(file_path: str) -> str:
    """
    Converts a .pptx file to Markdown.
    Extracts text from all slides.
    """
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)
    except Exception as e:
        return f"Error converting PPTX file: {e}"
