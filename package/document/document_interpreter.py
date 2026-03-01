"""
高级文档解释器工具
支持多种文档格式（PDF, Word, Excel, PPTX, EPUB, Markdown等）的解析、数据分析和AI摘要提问。
"""
import os
import sys
import shutil
import re
import csv
import json
import argparse
from collections import Counter
from typing import Optional, Dict, List, Any

# Third-party imports
import docx
from pypdf import PdfReader, PdfWriter
import pdfplumber
import pandas as pd
from openpyxl import load_workbook, Workbook
from pptx import Presentation
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import requests
from dotenv import load_dotenv

# Local imports
# Use consistent project root resolution
project_root = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
markitdown_src = os.path.join(project_root, 'markitdown', 'src')
if markitdown_src not in sys.path:
    sys.path.insert(0, markitdown_src)

try:
    from markitdown import MarkItDown
    MARKITDOWN_AVAILABLE = True
except ImportError:
    # Try importing from installed package if not in local src
    try:
        from markitdown import MarkItDown
        MARKITDOWN_AVAILABLE = True
    except ImportError:
        MARKITDOWN_AVAILABLE = False

from package.core_utils.log_manager import LogManager
from package.core_utils.config_loader import config_loader

logger = LogManager.get_logger(__name__)

class DocumentInterpreter:
    """
    一个高级文档解释器，支持多种格式的读取、分析和AI辅助处理。
    """
    def __init__(self, api_key: Optional[str] = None):
        self.api_key = api_key or config_loader.get("api.deepseek.key")
        try:
            self.mid = MarkItDown() if MARKITDOWN_AVAILABLE else None
        except Exception as e:
            logger.warning(f"Failed to initialize MarkItDown: {e}")
            self.mid = None
        self.deepseek_url = config_loader.get("api.deepseek.endpoint", "https://api.deepseek.com/v1") + "/chat/completions"

    def get_file_type(self, file_path: str) -> Optional[str]:
        _, ext = os.path.splitext(file_path)
        return ext.lower()

    def read_text(self, file_path: str) -> str:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def read_pdf(self, file_path: str, extract_tables: bool = False) -> str:
        text = ""
        try:
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                    if extract_tables:
                        tables = page.extract_tables()
                        for table in tables:
                            text += "\n[Table]\n"
                            for row in table:
                                text += " | ".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            if not text.strip():
                # Fallback to pypdf if pdfplumber fails or returns empty
                reader = PdfReader(file_path)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            logger.error(f"Error reading PDF {file_path}: {e}")
            return f"Error reading PDF: {e}"
        return text

    def read_docx(self, file_path: str) -> str:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])

    def read_pptx(self, file_path: str) -> str:
        pres = Presentation(file_path)
        text_runs = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    def read_xlsx(self, file_path: str) -> str:
        df = pd.read_excel(file_path)
        return df.to_string()

    def read_csv(self, file_path: str) -> str:
        df = pd.read_csv(file_path)
        return df.to_string()

    def read_epub(self, file_path: str) -> str:
        book = epub.read_epub(file_path)
        text = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text.append(soup.get_text())
        return "\n".join(text)

    def read_markdown(self, file_path: str) -> str:
        return self.read_text(file_path)

    def convert_via_markitdown(self, file_path: str) -> str:
        if self.mid:
            try:
                result = self.mid.convert(file_path)
                return result.text_content
            except Exception as e:
                logger.error(f"MarkItDown conversion failed: {e}")
                return f"MarkItDown error: {e}"
        return "MarkItDown not available."

    def interpret(self, file_path: str) -> str:
        ext = self.get_file_type(file_path)
        logger.info(f"Interpreting file: {file_path} (Type: {ext})")

        if ext == '.txt':
            return self.read_text(file_path)
        elif ext == '.pdf':
            return self.read_pdf(file_path, extract_tables=True)
        elif ext in ['.docx', '.doc']:
            return self.read_docx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return self.read_xlsx(file_path)
        elif ext == '.csv':
            return self.read_csv(file_path)
        elif ext == '.pptx':
            return self.read_pptx(file_path)
        elif ext == '.epub':
            return self.read_epub(file_path)
        elif ext == '.md':
            return self.read_markdown(file_path)
        else:
            # Try markitdown for everything else
            if MARKITDOWN_AVAILABLE:
                return self.convert_via_markitdown(file_path)
            return f"Unsupported file type: {ext}"

    def count_word_frequency(self, text: str) -> Counter:
        words = re.findall(r'\b\w+\b', text.lower())
        return Counter(words)

    def extract_keywords(self, text: str, num_keywords: int = 5) -> List[tuple]:
        words = re.findall(r'\b\w+\b', text.lower())
        # Filter out common stop words (simplified)
        stop_words = {'the', 'a', 'an', 'and', 'or', 'but', 'is', 'are', 'was', 'were', 'to', 'of', 'in', 'on', 'with'}
        filtered_words = [w for w in words if w not in stop_words and len(w) > 1]
        word_counts = Counter(filtered_words)
        return word_counts.most_common(num_keywords)

    def analyze_data(self, file_path: str) -> str:
        ext = self.get_file_type(file_path)
        if ext in ['.csv', '.xlsx', '.xls']:
            try:
                df = pd.read_csv(file_path) if ext == '.csv' else pd.read_excel(file_path)
                # Select all columns to describe, including numeric and categorical
                summary = df.describe(include='all').to_string()
                return f"Data Summary:\n{summary}"
            except Exception as e:
                return f"Error analyzing data: {e}"
        return "Not a tabular data file."

    def summarize(self, text: str) -> str:
        if not self.api_key:
            return "DeepSeek API key not found. Cannot summarize."

        logger.info("Generating summary via DeepSeek")
        payload = {
            "model": "deepseek-chat",
            "messages": [
                {"role": "system", "content": "You are a professional document analyst. Provide a concise summary of the following text in Chinese."},
                {"role": "user", "content": f"Please summarize this document:\n\n{text[:10000]}"} # Limit to first 10k chars
            ],
            "max_tokens": 1000
        }
        try:
            headers = {"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"}
            response = requests.post(self.deepseek_url, headers=headers, json=payload, timeout=30)
            response.raise_for_status()
            return response.json()['choices'][0]['message']['content']
        except Exception as e:
            logger.error(f"Summarization failed: {e}")
            return f"Summarization failed: {e}"

    def detect_language(self, text: str) -> str:
        # Simple heuristic for language detection
        chinese_chars = re.findall(r'[\u4e00-\u9fff]', text[:1000])
        if len(chinese_chars) > len(text[:1000]) * 0.1:
            return "zh"
        return "en"

    def ask_question(self, text: str, question: str) -> str:
        if not self.api_key:
            return "DeepSeek API key not found."

        payload = {
            "model": "deepseek-chat",
            "messages": [
                {"role": "system", "content": "You are a helpful assistant. Use the provided document content to answer the user's question in Chinese."},
                {"role": "user", "content": f"Document content:\n\n{text[:15000]}\n\nQuestion: {question}"}
            ]
        }
        try:
            headers = {"Authorization": f"Bearer {self.api_key}", "Content-Type": "application/json"}
            response = requests.post(self.deepseek_url, headers=headers, json=payload, timeout=60)
            response.raise_for_status()
            return response.json()['choices'][0]['message']['content']
        except Exception as e:
            return f"Failed to answer question: {e}"

    # --- Editing Functions (Restored and Improved) ---

    def search_and_replace(self, text: str, search_term: str, replace_term: str) -> str:
        return text.replace(search_term, replace_term)

    def extract_specific_pages(self, input_file: str, output_file: str, page_numbers: List[int]):
        try:
            reader = PdfReader(input_file)
            writer = PdfWriter()
            for page_num in page_numbers:
                if 0 < page_num <= len(reader.pages):
                    writer.add_page(reader.pages[page_num - 1])
            with open(output_file, 'wb') as out_f:
                writer.write(out_f)
            return f"Successfully extracted pages to {output_file}"
        except Exception as e:
            return f"Error extracting PDF pages: {e}"

    def update_docx(self, file_path: str, new_content: str, paragraph_index: int = 0):
        try:
            backup_file = file_path + '.bak'
            shutil.copyfile(file_path, backup_file)
            doc = docx.Document(file_path)
            if paragraph_index < len(doc.paragraphs):
                doc.paragraphs[paragraph_index].text = new_content
                doc.save(file_path)
                return f"Successfully updated {file_path}"
            else:
                return "Paragraph index out of range."
        except Exception as e:
            return f"Error updating DOCX: {e}"

    def clean_csv(self, file_path: str):
        try:
            df = pd.read_csv(file_path)
            df.dropna(how='all', inplace=True) # Remove empty rows
            df = df.apply(lambda x: x.str.strip() if x.dtype == "object" else x) # Strip strings
            df.to_csv(file_path, index=False)
            return f"Successfully cleaned {file_path}"
        except Exception as e:
            return f"Error cleaning CSV: {e}"

    def export_to_excel(self, data: List[List[Any]], output_file: str):
        try:
            df = pd.DataFrame(data)
            df.to_excel(output_file, index=False, header=False)
            return f"Successfully exported to {output_file}"
        except Exception as e:
            return f"Error exporting to Excel: {e}"

def run(*args, **kwargs):
    """Entry point for ExtensionManager."""
    interpreter = DocumentInterpreter()
    main()

def main():
    parser = argparse.ArgumentParser(description="Document Interpreter Tool")
    parser.add_argument("file", nargs="?", help="Path to the document file")
    parser.add_argument("--summarize", action="store_true", help="Summarize the document")
    parser.add_argument("--analyze", action="store_true", help="Perform data analysis (CSV/Excel)")
    args = parser.parse_args()

    interpreter = DocumentInterpreter()

    if args.file:
        file_path = args.file
    else:
        file_path = input("请输入文件路径: ").strip('"').strip("'")

    if not os.path.exists(file_path):
        print(f"文件不存在: {file_path}")
        return

    content = interpreter.interpret(file_path)

    if args.summarize:
        print("\n--- 摘要 ---")
        print(interpreter.summarize(content))
        return

    if args.analyze:
        print("\n--- 数据分析 ---")
        print(interpreter.analyze_data(file_path))
        return

    print(f"\n检测到文件类型: {interpreter.get_file_type(file_path)}")
    print(f"检测到语言: {interpreter.detect_language(content)}")

    while True:
        print("\n选择操作:")
        print("1. 查看全文 (前2000字)")
        print("2. 词频统计")
        print("3. 关键词提取")
        print("4. AI 自动摘要")
        print("5. 数据分析 (仅限 CSV/Excel)")
        print("6. 提问 (针对文档内容)")
        print("7. 更多工具 (PDF提取/Word更新等)")
        print("8. 退出")

        choice = input("请输入选项 (1-8): ")

        if choice == '1':
            print(f"\n--- 内容预览 ---\n{content[:2000]}...")
        elif choice == '2':
            freq = interpreter.count_word_frequency(content)
            print(f"\n--- 词频统计 (Top 20) ---\n{freq.most_common(20)}")
        elif choice == '3':
            num = int(input("提取多少个关键词? ") or 5)
            keywords = interpreter.extract_keywords(content, num)
            print(f"\n--- 关键词 ---\n{keywords}")
        elif choice == '4':
            print("\n正在生成摘要...")
            print(interpreter.summarize(content))
        elif choice == '5':
            print(f"\n--- 数据分析 ---\n{interpreter.analyze_data(file_path)}")
        elif choice == '6':
            q = input("你想问关于文档的什么问题? ")
            print("\n正在思考...")
            print(interpreter.ask_question(content, q))
        elif choice == '7':
            print("\n--- 更多工具 ---")
            print("a. 搜索并替换 (仅限预览模式)")
            print("b. 提取 PDF 指定页面")
            print("c. 更新 Word 文档首段")
            print("d. 清洗 CSV 文件")
            sub_choice = input("请选择 (a/b/c/d): ")
            if sub_choice == 'a':
                s = input("搜索词: ")
                r = input("替换词: ")
                print(f"结果预览:\n{interpreter.search_and_replace(content[:1000], s, r)}")
            elif sub_choice == 'b' and interpreter.get_file_type(file_path) == '.pdf':
                pages = input("请输入页码 (逗号分隔, 如 1,3,5): ")
                page_list = [int(p.strip()) for p in pages.split(',')]
                out = input("输出文件名: ")
                print(interpreter.extract_specific_pages(file_path, out, page_list))
            elif sub_choice == 'c' and interpreter.get_file_type(file_path) == '.docx':
                text = input("请输入新内容: ")
                print(interpreter.update_docx(file_path, text))
            elif sub_choice == 'd' and interpreter.get_file_type(file_path) == '.csv':
                print(interpreter.clean_csv(file_path))
            else:
                print("无效选择或文件类型不匹配。")
        elif choice == '8':
            break
        else:
            print("无效选项。")

if __name__ == '__main__':
    main()
