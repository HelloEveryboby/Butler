"""
Marker-Lite 工具: 快速且准确地将文档转换为 Markdown、JSON、区块和 HTML。

主要功能:
1. 转换所有语言的 PDF、图片、PPTX、DOCX、XLSX、HTML、EPUB 文件。
2. 精准提取表格、方程、内联数学、链接、参考和代码块。
3. 自动提取并保存文档中的图像。
4. 去除页眉、页脚及其他文档瑕疵。
5. 支持通过 JSON schema 进行结构化提取（Beta）。
6. 使用 DeepSeek API 增强转换质量和结构化能力。

用法:
- 作为工具运行: python -m package.marker_tool <文件路径> --format <markdown|json|chunks|html>
- 结构化提取: python -m package.marker_tool <文件路径> --schema <schema.json>
"""

import os
import re
import json
import base64
import argparse
from typing import List, Dict, Any, Optional, Union
from pathlib import Path

# Local Parsing Libraries
import pdfplumber
import pandas as pd
from docx import Document
from pptx import Presentation
from bs4 import BeautifulSoup
from ebooklib import epub
import ebooklib
from PIL import Image
import requests
from dotenv import load_dotenv

try:
    import pytesseract
    PYTESSERACT_AVAILABLE = True
except ImportError:
    PYTESSERACT_AVAILABLE = False

try:
    from aip import AipOcr
    BAIDU_OCR_AVAILABLE = True
except ImportError:
    BAIDU_OCR_AVAILABLE = False

from package.core_utils.log_manager import LogManager
from package.core_utils.config_loader import config_loader

logger = LogManager.get_logger(__name__)

class MarkerTool:
    """
    一个轻量级的 Marker 实现，使用 DeepSeek API 进行后端增强。
    """
    def __init__(self, api_key: Optional[str] = None, base_url: str = None):
        self.api_key = api_key or config_loader.get("api.deepseek.key")
        self.base_url = base_url or config_loader.get("api.deepseek.endpoint", "https://api.deepseek.com")
        self._init_baidu_ocr()

    def _init_baidu_ocr(self):
        self.baidu_ocr = None
        if BAIDU_OCR_AVAILABLE:
            app_id = config_loader.get("api.baidu.app_id")
            api_key = config_loader.get("api.baidu.api_key")
            secret_key = config_loader.get("api.baidu.secret_key")
            if app_id and api_key and secret_key:
                try:
                    self.baidu_ocr = AipOcr(app_id, api_key, secret_key)
                except Exception as e:
                    logger.warning(f"Failed to init Baidu OCR: {e}")
        self.output_dir = Path("data/marker_output")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.image_dir = self.output_dir / "images"
        self.image_dir.mkdir(exist_ok=True)

    def _get_deepseek_response(self, prompt: str, system_prompt: str = "You are a professional document converter.", json_mode: bool = False) -> str:
        if not self.api_key:
            return "Error: DeepSeek API Key not found in config or env."

        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "Content-Type": "application/json"
        }
        payload = {
            "model": "deepseek-chat",
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.1
        }
        if json_mode:
            payload["response_format"] = {"type": "json_object"}

        try:
            response = requests.post(f"{self.base_url}/v1/chat/completions", headers=headers, json=payload, timeout=120)
            response.raise_for_status()
            return response.json()['choices'][0]['message']['content']
        except Exception as e:
            logger.error(f"DeepSeek API error: {e}")
            return f"Error: {e}"

    def extract_pdf(self, file_path: str) -> Dict[str, Any]:
        """使用 pdfplumber 提取文本和图像"""
        raw_content = []
        images = []

        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages):
                # 提取文本
                text = page.extract_text() or ""
                raw_content.append(f"--- Page {i+1} ---\n{text}")

                # 提取图像
                for j, img in enumerate(page.images):
                    try:
                        # 简单的图像保存逻辑
                        img_obj = page.within_bbox((img["x0"], img["top"], img["x1"], img["bottom"])).to_image()
                        img_name = f"page_{i+1}_img_{j+1}.png"
                        img_path = self.image_dir / img_name
                        img_obj.save(img_path)
                        images.append(str(img_path))
                    except Exception as e:
                        logger.warning(f"Failed to extract image on page {i+1}: {e}")

        return {"text": "\n\n".join(raw_content), "images": images}

    def extract_docx(self, file_path: str) -> Dict[str, Any]:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        images = []

        # 尝试提取 DOCX 中的图像
        try:
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    img_name = f"docx_img_{len(images)+1}.png"
                    img_path = self.image_dir / img_name
                    with open(img_path, "wb") as f:
                        f.write(rel.target_part.blob)
                    images.append(str(img_path))
        except Exception as e:
            logger.warning(f"Failed to extract images from DOCX: {e}")

        return {"text": text, "images": images}

    def extract_pptx(self, file_path: str) -> Dict[str, Any]:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return {"text": "\n".join(text), "images": []}

    def extract_xlsx(self, file_path: str) -> Dict[str, Any]:
        df = pd.read_excel(file_path)
        return {"text": df.to_markdown(), "images": []}

    def extract_epub(self, file_path: str) -> Dict[str, Any]:
        book = epub.read_epub(file_path)
        text = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text.append(soup.get_text())
        return {"text": "\n".join(text), "images": []}

    def split_into_chunks(self, text: str, max_chunk_size: int = 1000) -> List[str]:
        """简单的语义块切分逻辑"""
        # 优先按二级标题切分
        sections = re.split(r'\n(?=## )', text)
        chunks = []
        for section in sections:
            if len(section) > max_chunk_size:
                # 如果依然过大，按段落切分
                paragraphs = section.split('\n\n')
                current_chunk = ""
                for p in paragraphs:
                    if len(current_chunk) + len(p) < max_chunk_size:
                        current_chunk += p + "\n\n"
                    else:
                        if current_chunk: chunks.append(current_chunk.strip())
                        current_chunk = p + "\n\n"
                if current_chunk: chunks.append(current_chunk.strip())
            else:
                chunks.append(section.strip())
        return chunks

    def convert(self, file_path: str, output_format: str = "markdown", json_schema: Optional[Dict] = None, custom_prompt: Optional[str] = None, skip_confirmation: bool = True) -> Union[str, List[str], Dict]:
        """
        skip_confirmation defaults to True here because we move the confirmation logic
        to the intent handler to avoid blocking background threads.
        """
        ext = Path(file_path).suffix.lower()
        logger.info(f"Converting {file_path} to {output_format}")

        # 1. 本地提取
        if ext == '.pdf':
            extracted = self.extract_pdf(file_path)
        elif ext == '.docx':
            extracted = self.extract_docx(file_path)
        elif ext == '.pptx':
            extracted = self.extract_pptx(file_path)
        elif ext in ['.xlsx', '.xls']:
            extracted = self.extract_xlsx(file_path)
        elif ext == '.epub':
            extracted = self.extract_epub(file_path)
        elif ext in ['.html', '.htm']:
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                extracted = {"text": soup.get_text(), "images": []}
        elif ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                extracted = {"text": f.read(), "images": []}
        elif ext in ['.png', '.jpg', '.jpeg']:
            text = ""
            # Priority 1: Baidu OCR (Online, accurate)
            if self.baidu_ocr:
                try:
                    with open(file_path, 'rb') as f:
                        image = f.read()
                    res = self.baidu_ocr.basicGeneral(image)
                    if "words_result" in res:
                        text = "\n".join([w["words"] for w in res["words_result"]])
                    else:
                        logger.warning(f"Baidu OCR error: {res}")
                except Exception as e:
                    logger.warning(f"Baidu OCR failed: {e}")

            # Priority 2: Tesseract (Offline)
            if not text and PYTESSERACT_AVAILABLE:
                try:
                    text = pytesseract.image_to_string(Image.open(file_path), lang='chi_sim+eng')
                except Exception as e:
                    logger.warning(f"Tesseract OCR failed: {e}")

            if not text:
                text = f"[OCR not available or failed. Image File: {file_path}]"

            extracted = {"text": text, "images": [file_path]}
        else:
            return f"Unsupported extension: {ext}"

        # 2. LLM 增强与转换
        raw_text = extracted["text"]

        if not skip_confirmation:
            # Note: This blocking input is discouraged in background threads.
            # It's kept for CLI direct usage, but intent handlers should pass skip_confirmation=True.
            print(f"\n--- 预解析完成 ---")
            print(f"文件: {file_path}")
            print(f"提取文本长度: {len(raw_text)} 字符")
            print(f"提取图像数量: {len(extracted.get('images', []))}")
            try:
                confirm = input("是否继续调用 DeepSeek API 进行精准转换/提取? (y/n): ")
                if confirm.lower() != 'y':
                    return "用户取消转换。"
            except EOFError:
                logger.warning("Standard input not available for confirmation.")

        if json_schema:
            system_prompt = "You are a specialized data extractor. Extract information from the document text and format it strictly according to the provided JSON schema."
            prompt = f"Document Text:\n{raw_text[:20000]}\n\nJSON Schema:\n{json.dumps(json_schema, indent=2)}\n\nPlease provide the extracted JSON."
            return self._get_deepseek_response(prompt, system_prompt, json_mode=True)

        # 常规转换模式
        system_prompt = (
            "You are an expert document converter similar to 'Marker'. "
            "Your goal is to convert document text into high-quality Markdown. "
            "Instructions:\n"
            "1. Remove headers, footers, page numbers, and repetitive artifacts.\n"
            "2. Properly format tables using Markdown syntax.\n"
            "3. Convert math equations to LaTeX format (e.g., $...$ or $$...$$).\n"
            "4. Preserve the logical structure (headings, lists, code blocks).\n"
            "5. If images were extracted, I will provide their markers; leave placeholders if needed.\n"
            "6. Output ONLY the converted content in the requested format."
        )

        if custom_prompt:
            prompt = f"{custom_prompt}\n\nDocument Text:\n{raw_text[:30000]}"
        else:
            format_instruction = f"Convert the following text to {output_format.upper() if output_format != 'chunks' else 'MARKDOWN'}:"
            prompt = f"{format_instruction}\n\n{raw_text[:30000]}"

        converted_text = self._get_deepseek_response(prompt, system_prompt)

        if output_format == "chunks":
            return self.split_into_chunks(converted_text)

        return converted_text

def run(*args, **kwargs):
    """Extension Manager entry point."""
    tool = MarkerTool()

    # CLI Argument Parsing
    parser = argparse.ArgumentParser(description="Marker-Lite Document Converter")
    parser.add_argument("file", help="Path to the document")
    parser.add_argument("-f", "--format", default="markdown", choices=["markdown", "json", "html", "chunks"], help="Output format")
    parser.add_argument("-s", "--schema", help="Path to JSON schema file for structured extraction")
    parser.add_argument("-p", "--prompt", help="Custom prompt for conversion")
    parser.add_argument("-y", "--yes", action="store_true", help="Skip confirmation")

    # In Jarvis, args might be passed as a list
    if args and isinstance(args[0], list):
        parsed_args = parser.parse_args(args[0])
    else:
        # Fallback if called with kwargs or other means
        file_path = kwargs.get("file")
        if not file_path:
             print("Error: No file path provided.")
             return
        tool_format = kwargs.get("format", "markdown")
        schema_path = kwargs.get("schema")
        prompt_text = kwargs.get("prompt")
        skip_confirm = kwargs.get("yes", False)

        # Mocking parser namespace for consistency
        class Args: pass
        parsed_args = Args()
        parsed_args.file = file_path
        parsed_args.format = tool_format
        parsed_args.schema = schema_path
        parsed_args.prompt = prompt_text
        parsed_args.yes = skip_confirm

    if not os.path.exists(parsed_args.file):
        print(f"File not found: {parsed_args.file}")
        return

    schema = None
    if parsed_args.schema and os.path.exists(parsed_args.schema):
        with open(parsed_args.schema, 'r', encoding='utf-8') as f:
            schema = json.load(f)

    result = tool.convert(parsed_args.file, parsed_args.format, schema, parsed_args.prompt, parsed_args.yes)
    print(result)

if __name__ == "__main__":
    import sys
    run(sys.argv[1:])
